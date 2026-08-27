"""生成答辩 PPT（python-pptx）。

运行：python deliverables/build_ppt.py
输出：deliverables/答辩PPT.pptx
设计：Midnight Executive 配色（navy/ice/white + gold accent），深色封面/结尾、浅色内容页，
统一「编号圆徽」视觉母题，结果页使用大数字 stat callout 与对比条。
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

NAVY = RGBColor(0x1E, 0x27, 0x61)
ICE = RGBColor(0xCA, 0xDC, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD = RGBColor(0xF9, 0xA8, 0x26)
DARKTEXT = RGBColor(0x22, 0x2A, 0x44)
GREY = RGBColor(0x5A, 0x60, 0x72)
LIGHTBG = RGBColor(0xF4, 0xF7, 0xFB)

EMU = 914400
SW, SH = 13.333, 7.5

prs = Presentation()
prs.slide_width = Emu(int(SW * EMU))
prs.slide_height = Emu(int(SH * EMU))
BLANK = prs.slide_layouts[6]


def _in(v: float) -> Emu:
    return Emu(int(v * EMU))


def slide():
    return prs.slides.add_slide(BLANK)


def bg(s, color):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color


def rect(s, x, y, w, h, color, shape=MSO_SHAPE.RECTANGLE, line=None):
    sp = s.shapes.add_shape(shape, _in(x), _in(y), _in(w), _in(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(_in(x), _in(y), _in(w), _in(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, (txt, size, color, bold, *rest) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        font_name = rest[0] if rest else None
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        if font_name:
            r.font.name = font_name
    return tb


def badge(s, x, y, n: str):
    c = rect(s, x, y, 0.55, 0.55, GOLD, MSO_SHAPE.OVAL)
    tf = c.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = n
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = NAVY


def content_header(s, num: str, title: str):
    bg(s, LIGHTBG)
    rect(s, 0, 0, SW, 1.25, NAVY)
    badge(s, 0.6, 0.35, num)
    text(s, 1.4, 0.32, 11.0, 0.7, [(title, 30, WHITE, True, "Calibri")], anchor=MSO_ANCHOR.MIDDLE)


def bullets(s, x, y, w, items, size=17, gap=True):
    tb = s.shapes.add_textbox(_in(x), _in(y), _in(w), _in(5.2))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    for i, (head, desc) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10 if gap else 4)
        r = p.add_run()
        r.text = f"▸ {head}  "
        r.font.size = Pt(size)
        r.font.bold = True
        r.font.color.rgb = NAVY
        if desc:
            r2 = p.add_run()
            r2.text = desc
            r2.font.size = Pt(size - 2)
            r2.font.color.rgb = GREY
    return tb


def card(s, x, y, w, h, title, body, accent=GOLD):
    rect(s, x, y, w, h, WHITE)
    rect(s, x, y, 0.12, h, accent)
    text(s, x + 0.3, y + 0.22, w - 0.45, 0.5, [(title, 16, NAVY, True)])
    text(s, x + 0.3, y + 0.78, w - 0.45, h - 0.9, [(body, 12.5, GREY, False)])


def stat(s, x, y, w, big, label, color=NAVY):
    rect(s, x, y, w, 1.9, WHITE)
    text(s, x, y + 0.25, w, 0.95, [(big, 46, color, True)], align=PP_ALIGN.CENTER)
    text(s, x, y + 1.25, w, 0.5, [(label, 13, GREY, False)], align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------- 1 封面
s = slide()
bg(s, NAVY)
rect(s, 0, 3.05, SW, 0.06, GOLD)
text(s, 1.0, 1.7, 11.3, 1.4,
     [("面向多智能体协作的低开销通信、", 38, WHITE, True, "Georgia"),
      ("状态传递与共享记忆机制", 38, WHITE, True, "Georgia")], anchor=MSO_ANCHOR.MIDDLE)
text(s, 1.0, 3.3, 11.3, 0.6, [("OS 大赛 · 赛题 10 · 应用创新（社区赛题）", 20, ICE, False)])
text(s, 1.0, 4.1, 11.3, 0.6,
     [("系统层机制：结构化通信 · 非文本状态 · 共享记忆 · CodeAct 沙箱", 16, GOLD, True)])
text(s, 1.0, 6.4, 11.3, 0.5, [("目标环境：openEuler 24.03-LTS-SP3 (x86_64)", 14, ICE, False)])

# ---------------------------------------------------------------- 2 问题与动机
s = slide()
content_header(s, "1", "问题与动机")
text(s, 0.6, 1.6, 12.1, 0.6,
     [("当前多 Agent 协作以自然语言 / JSON 透传中间结果，三大痛点：", 18, DARKTEXT, True)])
card(s, 0.6, 2.4, 3.95, 3.2, "① 通信冗长", "重复上下文多，token 消耗高，协作成本随轮次与 Agent 数放大。")
card(s, 4.75, 2.4, 3.95, 3.2, "② 反复编解码", "中间状态在「内部状态—文本—内部状态」间转换，增加时延与语义损耗。", ICE)
card(s, 8.9, 2.4, 3.85, 3.2, "③ 经验难复用", "中间知识与经验难以沉淀，相似任务仍从头开始，缺乏持续积累。")
text(s, 0.6, 5.9, 12.1, 0.6,
     [("我们从「系统层机制」解决，而非普通工作流编排（LangGraph/AutoGen/CrewAI）。", 15, NAVY, True)])

# ---------------------------------------------------------------- 3 目标与评分
s = slide()
content_header(s, "2", "目标与评分点对照")
items = [
    ("通信效率 · 25", "结构化协议 + Artifact 引用 + 非文本状态，Token 节省"),
    ("状态传递创新 · 20", "embedding / compact vector / plan DAG 直传且被真实消费"),
    ("记忆复用效果 · 20", "跨任务共享记忆，Effective / Harmful 命中度量"),
    ("系统完整性 · 20", "多 Agent 运行时、≥10 轮稳定、故障恢复"),
    ("实验验证 · 15", "4 组主实验 + 消融，均值/std/P50/P95，原始数据"),
]
bullets(s, 0.8, 1.7, 11.6, items, size=19)

# ---------------------------------------------------------------- 4 总体架构
s = slide()
content_header(s, "3", "总体架构（13 模块 · 核心自研）")
mods = [
    ("Runtime", 0.6), ("Agents", 2.45), ("Protocol", 4.3), ("Registry", 6.15),
    ("Scheduler", 8.0), ("StateExch", 9.85),
]
for name, x in mods:
    rect(s, x, 1.7, 1.7, 0.9, NAVY)
    text(s, x, 1.7, 1.7, 0.9, [(name, 12, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
mods2 = [
    ("Artifact", 0.6), ("Memory", 2.45), ("Sandbox", 4.3), ("Observ.", 6.15),
    ("Evaluation", 8.0), ("Providers", 9.85),
]
for name, x in mods2:
    rect(s, x, 2.75, 1.7, 0.9, RGBColor(0x33, 0x40, 0x8C))
    text(s, x, 2.75, 1.7, 0.9, [(name, 12, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
card(s, 0.6, 4.1, 5.95, 2.6, "四类 Agent",
     "Planner（规划/DAG）· Retriever（检索/复用）· CodeAct-Executor（生成并执行代码）· Reviewer（审查/总结/沉淀记忆）。")
card(s, 6.8, 4.1, 5.95, 2.6, "设计原则",
     "核心 Runtime/协议/状态/记忆全部自研；引用而非搬运；非文本状态被真实消费；mock-first 可复现；全链路降级。", ICE)

# ---------------------------------------------------------------- 5 通信协议
s = slide()
content_header(s, "4", "结构化通信协议")
bullets(s, 0.8, 1.7, 6.0, [
    ("唯一真源", "AgentMessage（Pydantic 全字段）"),
    ("高密度语义", "动作 / 参数 / 结果 / 能力，替代自然语言"),
    ("协议机制", "握手 · 能力发现 · schema 校验 · 版本"),
    ("幂等 + trace", "相同请求复用结果；全链路 trace_id"),
    ("公平双模式", "text↔structured 同源映射"),
], size=17)
card(s, 7.1, 1.85, 5.6, 4.5, "引用而非搬运",
     "长内容（事实表 / 代码 / 日志 / 向量）只存一次，消息仅传：\n\n"
     "• 对象 ID / 内容哈希\n• artifact:// · state:// · memory:// · vector://\n• 摘要 + 访问元数据\n\n"
     "→ 通信效率 25 分的主要来源；引用失效自动降级文本。", GOLD)

# ---------------------------------------------------------------- 6 非文本状态
s = slide()
content_header(s, "5", "非文本状态传递（被接收方真实消费）")
card(s, 0.6, 1.7, 3.95, 2.3, "embedding / 语义向量", "由 embedding 生成，用于语义检索与路由。")
card(s, 4.75, 1.7, 3.95, 2.3, "compact state vector", "计划的紧凑数值特征，供接收方快速决策。", ICE)
card(s, 8.9, 1.7, 3.85, 2.3, "plan DAG", "结构化计划图，接收方据此还原 operands 并调度。")
text(s, 0.6, 4.3, 12.1, 0.6, [("生成 → 内容寻址存储 → state:// 引用 → 共享内存传输 → 接收方消费", 16, NAVY, True)])
text(s, 0.6, 5.05, 12.1, 1.6,
     [("Retriever 用 plan embedding 对检索目标做语义排序、用 DAG 还原 operands —— 不再重新解析文本。", 14.5, GREY, False),
      ("大向量经 multiprocessing.shared_memory 零拷贝传输（openEuler 可开启），不可用自动降级。", 14.5, GREY, False)])

# ---------------------------------------------------------------- 7 共享记忆
s = slide()
content_header(s, "6", "共享记忆复用")
bullets(s, 0.8, 1.7, 6.0, [
    ("统一记忆单元", "全字段 + 四类（Working/Episodic/Semantic/Procedural）"),
    ("混合检索", "语义 + 关键词 + 标签，× 质量权重"),
    ("质量控制", "去重 / 版本 / 错误记忆降权 / provenance"),
    ("命中分级", "Retrieved · Used · Effective · Harmful"),
    ("跨任务复用", "同一子问题复用结论；相似子问题复用策略"),
], size=17)
card(s, 7.1, 1.85, 5.6, 4.5, "强证据：连续任务 A / B",
     "A1 修复零除缺陷并沉淀「修复策略」；\nA2 修复相邻模块——不提供策略，\n仅靠复用 A1 的过程性记忆即可通过沙箱验收。\n\n"
     "无记忆时 A2 / B2 失败 → 证明带来\n跨任务可迁移能力，而非简单缓存。", GOLD)

# ---------------------------------------------------------------- 8 CodeAct 沙箱
s = slide()
content_header(s, "7", "CodeAct 与沙箱")
text(s, 0.6, 1.6, 12.1, 0.6, [("LLM 生成 Python 代码 → 受限环境执行 → 回传 stdout/stderr/exit/资源/artifact", 17, DARKTEXT, True)])
card(s, 0.6, 2.4, 3.95, 3.6, "进程隔离", "subprocess + timeout + 独立进程组超时清理；Python -I 隔离模式。")
card(s, 4.75, 2.4, 3.95, 3.6, "资源限制", "setrlimit：CPU / 内存 / 文件 / 句柄 / 进程数；输出大小限制。", ICE)
card(s, 8.9, 2.4, 3.85, 3.6, "可降级增强", "工作目录隔离 + 环境白名单 + 网络禁用；cgroup v2 / Podman 可选，缺失自动降级。")

# ---------------------------------------------------------------- 9 实验设计
s = slide()
content_header(s, "8", "实验设计（公平 · 可复现）")
bullets(s, 0.8, 1.7, 11.8, [
    ("四组主实验", "A 文本基线 / B 结构化 / C 结构化+状态 / D 全系统"),
    ("消融", "去记忆 / 去状态 / 去沙箱，隔离各机制贡献"),
    ("16 项指标 + 派生", "Token 节省率 / 时延改善率 / 重复计算降低率 / 有效记忆命中率"),
    ("公平性", "固定模型/温度/工具/重试/随机种子；mock-first 确定性"),
    ("可复现", "每配置 ×5，报告均值/std/P50/P95；结果带时间戳落 results/ 不覆盖"),
], size=18)

# ---------------------------------------------------------------- 10 实验结果
s = slide()
bg(s, NAVY)
text(s, 0.6, 0.5, 12.1, 0.8, [("实验结果（真实数据 · results/）", 30, WHITE, True, "Georgia")])
stat(s, 0.7, 1.7, 3.7, "85.1%", "Token 节省（C/D vs 文本基线）", NAVY)
stat(s, 4.8, 1.7, 3.7, "100%", "任务成功率（全部配置）", NAVY)
stat(s, 8.9, 1.7, 3.7, "1.00", "有效记忆命中率（无负迁移）", NAVY)
card(s, 0.7, 3.95, 5.85, 2.9, "消融定位贡献",
     "去非文本状态 → Token 节省由 85% 退回 6%（状态是节省主因）；\n"
     "去共享记忆 → 重复计算降低归零（复用是降本主因）。\n两机制正交、可叠加。")
card(s, 6.75, 3.95, 5.85, 2.9, "诚实的权衡",
     "CodeAct 沙箱引入子进程隔离，D 端到端时延高于基线；\n"
     "去沙箱配置时延即回落 → 安全/时延权衡，\n可用常驻 worker 摊薄（未来工作）。", ICE)

# ---------------------------------------------------------------- 11 稳定性与 openEuler
s = slide()
content_header(s, "9", "稳定性与 openEuler 适配")
card(s, 0.6, 1.7, 5.95, 4.8, "稳定性",
     "• ≥10 轮连续执行成功，RSS 增长受控\n\n"
     "• 故障注入与恢复：\n   - Agent 崩溃 → 失败隔离不崩溃\n   - LLM 暂时失败 → 重试恢复\n"
     "   - 记忆库不可用 → 降级\n   - 引用失效 → 降级文本")
card(s, 6.75, 1.7, 5.95, 4.8, "openEuler 24.03-LTS-SP3",
     "• dnf 一键部署，隔离 venv，幂等脚本\n\n"
     "• verify_openeuler.sh 输出：OS/内核/架构/\n   依赖/cgroup v2/共享内存/Socket/沙箱/\n   测试与 10 轮稳定性摘要\n\n"
     "• cgroup/共享内存/Podman 可用则启用，\n   不可用自动降级", ICE)

# ---------------------------------------------------------------- 12 总结
s = slide()
bg(s, NAVY)
rect(s, 0, 2.9, SW, 0.06, GOLD)
text(s, 1.0, 1.6, 11.3, 1.2,
     [("交付：多 Agent 协作的系统层基础设施", 34, WHITE, True, "Georgia")], anchor=MSO_ANCHOR.MIDDLE)
text(s, 1.0, 3.2, 11.3, 1.4,
     [("低开销 · 可解释 · 可复现 · 可在 openEuler 运行", 20, ICE, False)])
text(s, 1.0, 4.3, 11.3, 1.8,
     [("未来工作：常驻沙箱 worker · 真实 LLM 大规模评测 · 更强状态表示（KV/AST）· eBPF 观测 · 分布式运行时", 16, GOLD, True)])
text(s, 1.0, 6.5, 11.3, 0.5, [("谢谢！欢迎提问。", 18, WHITE, True)])


def main() -> None:
    out = Path(__file__).resolve().parent / "答辩PPT.pptx"
    prs.save(str(out))
    print(f"已生成: {out}（{len(prs.slides.__iter__.__self__._sldIdLst)} 页）")


if __name__ == "__main__":
    main()
